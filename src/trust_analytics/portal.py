"""Business-facing projection and export helpers for Trust Analytics Portal."""

from __future__ import annotations

import csv
import io
import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from trust_analytics.agents.quality_agent import QualityAgent
from trust_analytics.agents.sql_agent import SQLAgent
from trust_analytics.config import Settings, load_settings
from trust_analytics.data_loader import load_csvs
from trust_analytics.llm import make_client
from trust_analytics.metadata import case_root_from_data_dir, load_dbt_metadata
from trust_analytics.metrics import load_metrics_registry
from trust_analytics.models import (
    BusinessQuestion,
    PipelineItem,
    PipelineResult,
    QualityReport,
    ReviewMode,
    SQLAgentAnswer,
)
from trust_analytics.questions import REQUIRED_QUESTIONS, get_question, synthesize_business_question
from trust_analytics.workflow import run_pipeline, write_pipeline_outputs

CACHE_DIR = Path("outputs/sample")
ASK_RUNS_DIR = Path("outputs/ask")
HERO_ID = "q1_gtv_idr_by_asset_oct_2025"
AUDIT_ID = "q5_gtv_mom_trend_oct_dec_2025"
MTU_ID = "q3_mtu_oct_2025"

ASK_FIELD_LABELS = {
    "businessObjective": "Business objective",
    "period": "Time period",
    "segment": "Object or segment",
    "dimension": "Comparison dimension",
    "audience": "Audience",
    "desiredOutput": "Desired output",
}


@dataclass(frozen=True)
class CachedArtifacts:
    answers: list[SQLAgentAnswer]
    quality_reports: list[QualityReport]
    plans: list[dict[str, Any] | None]


def demo_questions() -> list[dict[str, str]]:
    copy = {
        HERO_ID: "Which asset class should we prioritise for next month's growth plan?",
        "q2_gtv_usd_oct_2025": "Can we use October GTV in USD for the leadership snapshot?",
        MTU_ID: "How many monthly transacting users did the platform have?",
        "q4_transaction_count_by_asset_oct_2025": "Which asset class drove the most October transaction activity?",
        AUDIT_ID: "Can we use the monthly summary for the quarter-end GTV trend?",
    }
    return [
        {
            "id": question.id,
            "text": copy.get(question.id, question.text),
            "metric": _metric_label(question.metric),
            "period": question.period,
        }
        for question in REQUIRED_QUESTIONS
    ]


def load_cached_artifacts(cache_dir: Path = CACHE_DIR) -> CachedArtifacts:
    answers = [
        SQLAgentAnswer.model_validate(item)
        for item in json.loads((cache_dir / "sql_agent_answers.json").read_text(encoding="utf-8"))
    ]
    quality = [
        QualityReport.model_validate(item)
        for item in json.loads((cache_dir / "quality_report.json").read_text(encoding="utf-8"))
    ]
    plans = json.loads((cache_dir / "question_plans.json").read_text(encoding="utf-8"))
    return CachedArtifacts(answers=answers, quality_reports=quality, plans=plans)


def _artifact_dirs() -> list[tuple[Path, str]]:
    dirs: list[tuple[Path, str]] = []
    if CACHE_DIR.exists():
        dirs.append((CACHE_DIR, "seed"))
    if ASK_RUNS_DIR.exists():
        for path in sorted(ASK_RUNS_DIR.iterdir()):
            if path.is_dir() and (path / "sql_agent_answers.json").exists():
                dirs.append((path, "ask"))
    return dirs


def _analysis_from_dir(
    analysis_id: str,
    artifact_dir: Path,
    *,
    source_kind: str,
    live_error: str | None = None,
) -> dict[str, Any] | None:
    artifacts = load_cached_artifacts(artifact_dir)
    for answer, quality, plan in zip(
        artifacts.answers, artifacts.quality_reports, artifacts.plans, strict=True
    ):
        if answer.question_id == analysis_id:
            analysis = project_analysis(
                answer=answer,
                quality=quality,
                question_plan=plan,
                from_cache=source_kind == "seed",
                live_error=live_error,
            )
            analysis["librarySource"] = source_kind
            analysis["artifactDir"] = str(artifact_dir)
            return analysis
    return None


def cached_analysis(analysis_id: str, *, live_error: str | None = None) -> dict[str, Any]:
    for artifact_dir, source_kind in _artifact_dirs():
        analysis = _analysis_from_dir(
            analysis_id,
            artifact_dir,
            source_kind=source_kind,
            live_error=live_error,
        )
        if analysis is not None:
            return analysis
    raise KeyError(f"No cached analysis found for {analysis_id}")


def list_cached_analyses() -> list[dict[str, Any]]:
    return [cached_analysis(q["id"]) for q in demo_questions()]


def list_library_items() -> list[dict[str, Any]]:
    items: list[dict[str, Any]] = []
    for artifact_dir, source_kind in _artifact_dirs():
        try:
            artifacts = load_cached_artifacts(artifact_dir)
        except FileNotFoundError:
            continue
        for answer, quality, plan in zip(
            artifacts.answers, artifacts.quality_reports, artifacts.plans, strict=True
        ):
            analysis = project_analysis(
                answer=answer,
                quality=quality,
                question_plan=plan,
                from_cache=source_kind == "seed",
            )
            analysis["librarySource"] = source_kind
            analysis["artifactDir"] = str(artifact_dir)
            items.append(
                {
                    "id": analysis["id"],
                    "source": source_kind,
                    "question": analysis["question"],
                    "metricName": analysis["metricName"],
                    "period": analysis["period"],
                    "status": analysis["status"],
                    "headline": analysis["headline"],
                    "decisionPack": analysis["decisionPack"],
                    "audit": analysis["audit"],
                    "artifactDir": str(artifact_dir),
                }
            )
    return items


def shape_question(
    *,
    question_text: str,
    business_objective: str | None = None,
    period: str | None = None,
    segment: str | None = None,
    dimension: str | None = None,
    audience: str | None = None,
    desired_output: str | None = None,
) -> dict[str, Any]:
    """Shape a messy business question into confirmable guided fields."""
    text = (question_text or "").strip()
    if not text:
        text = "Which asset class should we prioritise for next month's growth plan?"

    synthesized = synthesize_business_question(text, period_override=period or None)
    match = _match_verified_analysis(
        " ".join(
            item
            for item in [
                text,
                business_objective or "",
                period or "",
                segment or "",
                dimension or "",
                desired_output or "",
            ]
            if item
        )
    )
    analysis = cached_analysis(match["id"])
    inferred = _infer_ask_fields(
        text=text,
        analysis=analysis,
        synthesized_period=synthesized.period,
        business_objective=business_objective,
        period=period,
        segment=segment,
        dimension=dimension,
        audience=audience,
        desired_output=desired_output,
    )
    explicit = _explicit_ask_fields(
        text=text,
        synthesized_period=synthesized.period,
        business_objective=business_objective,
        period=period,
        segment=segment,
        dimension=dimension,
        audience=audience,
        desired_output=desired_output,
    )
    field_states = _field_states(
        inferred=inferred,
        explicit=explicit,
        provided={
            "businessObjective": business_objective,
            "period": period,
            "segment": segment,
            "dimension": dimension,
            "audience": audience,
            "desiredOutput": desired_output,
        },
    )
    missing = [
        {"key": key, "label": ASK_FIELD_LABELS[key]}
        for key in ASK_FIELD_LABELS
        if field_states[key]["status"] == "missing"
    ]
    ready_count = len(ASK_FIELD_LABELS) - len(missing)
    score = round(ready_count / len(ASK_FIELD_LABELS) * 100)
    blockers = [
        {
            "key": key,
            "label": ASK_FIELD_LABELS[key],
            "message": (
                f"Confirm {ASK_FIELD_LABELS[key].lower()} before validation, or make the "
                "question explicit enough for the agent to infer it."
            ),
        }
        for key, state in field_states.items()
        if state["status"] in {"missing", "inferred"}
    ]
    ready = not blockers
    return {
        "input": text,
        "canonicalQuestion": _canonical_question(analysis, inferred),
        "inferredMetric": analysis["metricName"],
        "inferredPeriod": inferred["period"],
        "inferredDimension": inferred["dimension"],
        "fields": inferred,
        "fieldStates": field_states,
        "confirmedFields": explicit,
        "missingFields": missing,
        "clarificationNeeds": blockers,
        "suggestedChips": _ask_suggested_chips(match["id"]),
        "recommendedAnalysisId": match["id"],
        "recommendedAnalysisTitle": analysis["decisionPack"]["title"],
        "verifiedPath": {
            "id": match["id"],
            "label": analysis["metricName"],
            "status": analysis["status"],
            "reason": match["reason"],
        },
        "confidence": match["confidence"],
        "quality": {
            "score": score,
            "label": "Ready to validate" if ready else "Needs confirmation",
            "ready": ready,
        },
        "explanation": (
            "The agent inferred the guided fields it can. Confirm or fill the "
            "remaining fields before running live validation."
        ),
    }


def run_live_pipeline(
    *,
    question_id: str | None = None,
    question_text: str | None = None,
    fields: dict[str, str] | None = None,
    settings: Settings | None = None,
) -> PipelineResult:
    settings = settings or load_settings()
    if not settings.db_path.exists():
        load_csvs(settings.data_dir, settings.db_path)

    question = _question_for_run(question_id=question_id, question_text=question_text, fields=fields)
    metadata = load_dbt_metadata(case_root_from_data_dir(settings.data_dir))
    registry = load_metrics_registry()
    llm = make_client(settings)
    sql_agent = SQLAgent(settings.db_path, metadata, llm, registry)
    quality_agent = QualityAgent(settings.db_path, registry, llm)
    return run_pipeline([question], sql_agent, quality_agent, ReviewMode.DEMO_APPROVE)


def run_live_analysis(
    *,
    question_id: str | None = None,
    question_text: str | None = None,
    settings: Settings | None = None,
) -> dict[str, Any]:
    result = run_live_pipeline(
        question_id=question_id,
        question_text=question_text,
        settings=settings,
    )
    item = result.items[0]
    return project_item(item, from_cache=False)


def run_and_persist_analysis(
    *,
    question_id: str | None = None,
    question_text: str | None = None,
    fields: dict[str, str] | None = None,
    output_root: Path = ASK_RUNS_DIR,
    settings: Settings | None = None,
) -> dict[str, Any]:
    result = run_live_pipeline(
        question_id=question_id,
        question_text=question_text,
        fields=fields,
        settings=settings,
    )
    item = result.items[0]
    analysis = project_item(item, from_cache=False)
    out = output_root / item.question.id
    write_pipeline_outputs(result, out)
    analysis["artifactDir"] = str(out)
    analysis["librarySource"] = "ask"
    return analysis


def run_or_cache(
    *, question_id: str | None = None, question_text: str | None = None
) -> dict[str, Any]:
    target_id = question_id or HERO_ID
    try:
        live = run_live_analysis(question_id=question_id, question_text=question_text)
        evidence = live.get("analystEvidence") or {}
        if not live.get("rows") or not evidence.get("sql"):
            return cached_analysis(
                target_id,
                live_error="Live run did not produce an executable decision result.",
            )
        return live
    except Exception as exc:  # noqa: BLE001 - public demo must degrade gracefully.
        return cached_analysis(target_id, live_error=f"{type(exc).__name__}: {exc}")


def _match_verified_analysis(text: str) -> dict[str, str]:
    lowered = text.lower()
    scores = {HERO_ID: 0, AUDIT_ID: 0, MTU_ID: 0}
    for token in ["asset", "asset class", "crypto", "gold", "growth", "priority", "prioritise", "prioritize"]:
        if token in lowered:
            scores[HERO_ID] += 2
    for token in ["trend", "mom", "month-on-month", "quarter", "december", "summary", "stale", "audit"]:
        if token in lowered:
            scores[AUDIT_ID] += 2
    for token in ["mtu", "user", "users", "customer", "transacting", "definition"]:
        if token in lowered:
            scores[MTU_ID] += 2
    if "gtv" in lowered:
        scores[HERO_ID] += 1
        scores[AUDIT_ID] += 1

    analysis_id = max(scores, key=scores.get)
    if scores[analysis_id] == 0:
        analysis_id = HERO_ID
    reason = {
        HERO_ID: "The question is closest to an asset-class growth decision.",
        AUDIT_ID: "The question is closest to a trend reliability and audit decision.",
        MTU_ID: "The question is closest to a metric-definition decision.",
    }[analysis_id]
    confidence = "high" if scores[analysis_id] >= 4 else "medium"
    return {"id": analysis_id, "reason": reason, "confidence": confidence}


def _metric_label(metric: str) -> str:
    labels = {
        "gtv_idr_by_asset_class": "GTV by asset class",
        "total_gtv_usd": "Total GTV in USD",
        "monthly_transacting_users": "Monthly transacting users",
        "transaction_count_by_asset_class": "Transaction count by asset class",
        "gtv_idr_month_on_month_trend": "GTV month-on-month trend",
    }
    return labels.get(metric, metric.replace("_", " ").title())


def _field_states(
    *,
    inferred: dict[str, str],
    explicit: dict[str, bool],
    provided: dict[str, str | None],
) -> dict[str, dict[str, str | bool]]:
    states: dict[str, dict[str, str | bool]] = {}
    for key, label in ASK_FIELD_LABELS.items():
        value = inferred.get(key) or ""
        user_value = (provided.get(key) or "").strip()
        if user_value and user_value != value:
            status = "manual_override"
        elif explicit.get(key):
            status = "confirmed"
        elif value:
            status = "inferred"
        else:
            status = "missing"
        states[key] = {
            "label": label,
            "value": value,
            "status": status,
            "confirmed": status in {"confirmed", "manual_override"},
        }
    return states


def _question_for_run(
    *,
    question_id: str | None,
    question_text: str | None,
    fields: dict[str, str] | None,
) -> BusinessQuestion:
    if question_id:
        return get_question(question_id)
    text = (question_text or "").strip()
    if not text:
        raise ValueError("Question text must be non-empty.")
    period = (fields or {}).get("period") or None
    return synthesize_business_question(text, period_override=period)


def _infer_ask_fields(
    *,
    text: str,
    analysis: dict[str, Any],
    synthesized_period: str,
    business_objective: str | None,
    period: str | None,
    segment: str | None,
    dimension: str | None,
    audience: str | None,
    desired_output: str | None,
) -> dict[str, str]:
    """Compute the effective value per field.

    A passed-in ``""`` means the user explicitly cleared the field, so we
    honor it and do not re-infer. ``None`` means the field was untouched by
    the user and the agent may infer from the question text or seed analysis.
    """
    lowered = text.lower()
    if business_objective is not None:
        objective = business_objective
    else:
        objective = ""
        if any(token in lowered for token in ["priority", "prioritise", "prioritize", "growth"]):
            objective = "Prioritise the next growth focus"
        elif any(token in lowered for token in ["trend", "safe", "rely", "audit"]):
            objective = "Decide whether this trend is safe to present"
        elif any(token in lowered for token in ["report", "mtu", "users"]):
            objective = "Choose the reporting definition"

    if period is not None:
        selected_period = period
    else:
        selected_period = (
            _display_period(synthesized_period) if synthesized_period != "unspecified" else ""
        )

    if segment is not None:
        selected_segment = segment
    else:
        selected_segment = ""
        if analysis["id"] == HERO_ID:
            selected_segment = "Completed trading activity"
        elif analysis["id"] == AUDIT_ID:
            selected_segment = "Completed GTV"
        elif analysis["id"] == MTU_ID:
            selected_segment = "Transacting users"

    if dimension is not None:
        selected_dimension = dimension
    else:
        selected_dimension = ""
        if analysis["id"] == HERO_ID:
            selected_dimension = "Asset class"
        elif analysis["id"] == AUDIT_ID:
            selected_dimension = "Month"

    return {
        "businessObjective": objective,
        "period": selected_period,
        "segment": selected_segment,
        "dimension": selected_dimension,
        "audience": audience if audience is not None else "",
        "desiredOutput": desired_output if desired_output is not None else "",
    }


def _explicit_ask_fields(
    *,
    text: str,
    synthesized_period: str,
    business_objective: str | None,
    period: str | None,
    segment: str | None,
    dimension: str | None,
    audience: str | None,
    desired_output: str | None,
) -> dict[str, bool]:
    """Return whether each field is explicitly committed by the user.

    A passed-in ``""`` means the user explicitly cleared the field, so we
    treat it as un-confirmed even if the original question text would have
    otherwise been a strong keyword hint.
    """
    lowered = text.lower()

    def _explicit(user_value: str | None, keywords: tuple[str, ...]) -> bool:
        if user_value == "":
            return False
        if user_value:
            return True
        return any(token in lowered for token in keywords)

    return {
        "businessObjective": _explicit(
            business_objective,
            ("priority", "prioritise", "prioritize", "growth", "report", "audit", "safe"),
        ),
        "period": (
            False
            if period == ""
            else bool(period or synthesized_period != "unspecified")
        ),
        "segment": _explicit(
            segment,
            ("completed", "trading", "users", "asset", "gtv"),
        ),
        "dimension": _explicit(
            dimension,
            ("asset class", "asset", "month", "trend", "definition", "source"),
        ),
        "audience": bool(audience),
        "desiredOutput": bool(desired_output),
    }


def _canonical_question(analysis: dict[str, Any], fields: dict[str, str]) -> str:
    period = fields["period"] or "the selected period"
    if analysis["id"] == HERO_ID:
        return (
            f"Which {fields['dimension'].lower()} should we prioritise in "
            f"{period} for the {fields['audience'].lower()} decision pack?"
        )
    if analysis["id"] == AUDIT_ID:
        return f"Can we use the {fields['period']} GTV trend for leadership reporting?"
    if analysis["id"] == MTU_ID:
        return f"Which MTU definition should {fields['audience'].lower()} use for {fields['period']}?"
    return analysis["question"]


def _display_period(period: str) -> str:
    labels = {
        "2025-10-01 to 2025-11-01": "October 2025",
        "2025-11-01 to 2025-12-01": "November 2025",
        "2025-12-01 to 2026-01-01": "December 2025",
    }
    return labels.get(period, period)


def _ask_suggested_chips(analysis_id: str) -> dict[str, list[str]]:
    common = {
        "audience": ["Leadership", "CEO", "Finance", "Growth team"],
        "desiredOutput": ["Decision pack", "Executive slide", "Email brief", "CSV extract"],
    }
    if analysis_id == AUDIT_ID:
        return {
            "businessObjective": ["Validate trend reliability", "Prepare audit brief"],
            "period": ["October to December 2025", "December 2025"],
            "segment": ["Completed GTV", "Monthly summary"],
            "dimension": ["Month", "Source"],
            **common,
        }
    if analysis_id == MTU_ID:
        return {
            "businessObjective": ["Choose reporting definition", "Report user activity"],
            "period": ["October 2025", "Last month"],
            "segment": ["Transacting users", "Product-active users"],
            "dimension": ["Definition", "Source"],
            **common,
        }
    return {
        "businessObjective": ["Prioritise growth focus", "Prepare leadership recommendation"],
        "period": ["October 2025", "Next month planning"],
        "segment": ["Completed trading activity", "All asset classes"],
        "dimension": ["Asset class", "Source"],
        **common,
    }


def project_item(item: PipelineItem, *, from_cache: bool) -> dict[str, Any]:
    plan = item.question_plan.model_dump(mode="json") if item.question_plan else None
    return project_analysis(
        answer=item.answer,
        quality=item.quality_report,
        question_plan=plan,
        from_cache=from_cache,
    )


def project_analysis(
    *,
    answer: SQLAgentAnswer,
    quality: QualityReport,
    question_plan: dict[str, Any] | None,
    from_cache: bool,
    live_error: str | None = None,
) -> dict[str, Any]:
    status = _business_status(answer, quality)
    rows = _normalise_rows(answer.result_rows)
    chart = _chart_spec(answer, rows)
    source_comparison = _source_comparison(quality)
    executive_summary = _executive_summary(answer, quality, status)
    use_this_for, do_not_use_for = _usage_boundaries(answer, status)
    return {
        "id": answer.question_id,
        "question": _business_question(answer),
        "metricName": _metric_name(answer),
        "period": answer.period,
        "status": status,
        "workflowSteps": _workflow_steps(answer, status),
        "headline": _headline(answer, rows),
        "decisionPrompt": _decision_prompt(answer),
        "recommendation": _recommendation(answer, status),
        "recommendedUse": _recommended_use(answer, status),
        "useThisFor": use_this_for,
        "doNotUseFor": do_not_use_for,
        "sourceCaveat": _source_caveat(answer, quality),
        "confidence": _confidence(answer, quality, status),
        "chartInsight": _chart_insight(answer, rows),
        "chart": chart,
        "rows": rows,
        "sourceComparison": source_comparison,
        "executiveSummary": executive_summary,
        "decisionPack": {
            "title": _pack_title(answer),
            "subtitle": "Leadership-ready brief generated from verified SQL evidence.",
            "deliverables": ["Executive slide", "CSV extract", "Email brief", "Evidence appendix"],
        },
        "emailDraft": make_email_draft(answer, executive_summary, status),
        "analystEvidence": {
            "sql": answer.sql,
            "source": answer.source.model_dump(mode="json") if answer.source else None,
            "derivationTrace": (
                answer.derivation_trace.model_dump(mode="json")
                if answer.derivation_trace
                else None
            ),
            "questionPlan": question_plan,
            "layerA": quality.layer_a.model_dump(mode="json"),
            "layerB": quality.layer_b.model_dump(mode="json"),
            "layerC": quality.layer_c.model_dump(mode="json"),
        },
        "audit": _audit_payload(answer, quality),
        "fromCache": from_cache,
        "liveError": live_error,
    }


def make_csv(analysis: dict[str, Any]) -> bytes:
    rows = analysis["rows"]
    if not rows:
        return b""
    out = io.StringIO()
    writer = csv.DictWriter(out, fieldnames=list(rows[0].keys()))
    writer.writeheader()
    writer.writerows(rows)
    return out.getvalue().encode("utf-8")


def make_email_draft(
    answer: SQLAgentAnswer, executive_summary: list[str], status: dict[str, str]
) -> dict[str, str]:
    subject = f"Decision pack ready: {_metric_name(answer)}"
    body = "\n".join(
        [
            "Hi,",
            "",
            "I prepared a leadership-ready decision pack in Trust Analytics.",
            "",
            *[f"- {line}" for line in executive_summary],
            "",
            f"Workflow status: {status['label']}.",
            "The pack includes the slide-ready recommendation, CSV extract, and analyst evidence appendix if challenged.",
            "",
            "Best,",
        ]
    )
    return {"subject": subject, "body": body}


def make_pptx(analysis: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_background(slide, RGBColor(248, 250, 252))
    _textbox(slide, Inches(0.62), Inches(0.32), Inches(3.6), Inches(0.32), "TRUST ANALYTICS", 10, RGBColor(71, 85, 105), bold=True)
    _textbox(slide, Inches(0.62), Inches(0.78), Inches(8.6), Inches(0.86), analysis["headline"], 27, RGBColor(15, 23, 42), bold=True)
    _pill(slide, Inches(10.3), Inches(0.76), analysis["status"]["label"])

    _textbox(slide, Inches(0.72), Inches(1.95), Inches(5.25), Inches(0.28), "DECISION RECOMMENDATION", 10, RGBColor(71, 85, 105), bold=True)
    _textbox(slide, Inches(0.72), Inches(2.32), Inches(5.45), Inches(1.04), analysis["recommendation"], 17, RGBColor(15, 23, 42), bold=True)
    _callout(slide, Inches(0.72), Inches(3.62), Inches(5.45), Inches(0.86), "Source caveat", analysis["sourceCaveat"], RGBColor(252, 245, 224), RGBColor(146, 93, 20))
    _boundary_box(slide, analysis, Inches(0.72), Inches(4.78), Inches(5.45), Inches(1.12))

    _textbox(slide, Inches(6.7), Inches(1.92), Inches(5.8), Inches(0.35), "Evidence-backed result", 13, RGBColor(15, 23, 42), bold=True)
    _draw_chart_panel(slide, analysis["chart"], Inches(6.7), Inches(2.38), Inches(5.65), Inches(3.34))
    _textbox(slide, Inches(6.7), Inches(5.92), Inches(5.6), Inches(0.42), analysis.get("chartInsight") or "", 11, RGBColor(71, 85, 105))
    _textbox(slide, Inches(0.72), Inches(6.55), Inches(11.8), Inches(0.28), "Generated from read-only SQL, source reconciliation, and analyst evidence. Appendix included for challenge.", 9, RGBColor(100, 116, 139))

    appendix = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_background(appendix, RGBColor(248, 250, 252))
    _textbox(appendix, Inches(0.62), Inches(0.45), Inches(8), Inches(0.45), "Appendix: challenge-ready evidence", 21, RGBColor(15, 23, 42), bold=True)
    evidence = analysis["analystEvidence"]
    qa = evidence["layerC"]["trust_profile"]
    source_rows = analysis.get("sourceComparison") or []
    source_note = "; ".join(
        f"{row['source']}: {row.get('deltaVsPrimary') if row.get('deltaVsPrimary') is not None else 'n/a'}%"
        for row in source_rows[:3]
    )
    body = [
        f"Primary source: {(evidence['source'] or {}).get('primary_table', 'unknown')}",
        f"Trust profile: {qa['overall']} | correctness={qa['dimensions']['correctness']} | source={qa['dimensions']['source_reliability']} | ambiguity={qa['dimensions']['ambiguity']}",
        f"Source reconciliation: {source_note or 'No comparison sources available.'}",
        f"QA summary: {qa['reviewer_summary']}",
        "",
        "SQL:",
        evidence["sql"][:1200],
    ]
    _textbox(appendix, Inches(0.75), Inches(1.18), Inches(11.8), Inches(5.7), "\n".join(body), 10.5, RGBColor(30, 41, 59))

    stream = io.BytesIO()
    prs.save(stream)
    return stream.getvalue()


def _textbox(slide, left, top, width, height, text, size, color, bold=False):
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _slide_background(slide, color):
    from pptx.util import Inches

    shape = slide.shapes.add_shape(1, 0, 0, Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def _callout(slide, left, top, width, height, label, text, fill, color):
    from pptx.dml.color import RGBColor

    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = fill
    _textbox(slide, left + width * 0.035, top + height * 0.12, width * 0.92, height * 0.18, label.upper(), 8.5, color, bold=True)
    _textbox(slide, left + width * 0.035, top + height * 0.36, width * 0.92, height * 0.46, text, 10.5, RGBColor(30, 41, 59))


def _boundary_box(slide, analysis, left, top, width, height):
    from pptx.dml.color import RGBColor

    _textbox(slide, left, top, width * 0.47, height * 0.18, "USE THIS FOR", 8.5, RGBColor(28, 113, 90), bold=True)
    _textbox(slide, left, top + height * 0.24, width * 0.47, height * 0.72, analysis.get("useThisFor", ""), 10.2, RGBColor(15, 23, 42))
    _textbox(slide, left + width * 0.53, top, width * 0.47, height * 0.18, "DO NOT USE FOR", 8.5, RGBColor(146, 64, 14), bold=True)
    _textbox(slide, left + width * 0.53, top + height * 0.24, width * 0.47, height * 0.72, analysis.get("doNotUseFor", ""), 10.2, RGBColor(15, 23, 42))


def _pill(slide, left, top, text):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(5, left, top, Inches(2.3), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(232, 244, 239)
    shape.line.color.rgb = RGBColor(62, 128, 104)
    tf = shape.text_frame
    tf.text = text
    tf.paragraphs[0].runs[0].font.size = Pt(11)
    tf.paragraphs[0].runs[0].font.bold = True
    tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(28, 97, 76)


def _draw_chart_panel(slide, chart, left, top, width, height):
    """Route the chart payload to the right renderer based on chart.type."""
    chart_type = chart.get("type") or "bar"
    values = chart.get("values") or []
    if chart_type == "table" or not values:
        _draw_table_panel(slide, chart, left, top, width, height)
    else:
        _draw_bar_chart(slide, chart, left, top, width, height)


def _draw_bar_chart(slide, chart, left, top, width, height):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    labels = chart.get("labels") or []
    values = chart.get("values") or []
    if not values:
        return
    paired = list(zip(labels, values, strict=False))
    max_value = max(values)
    rows = max(len(paired), 1)
    bar_h = height / rows * 0.45
    gap = height / rows * 0.55
    for idx, (label, value) in enumerate(paired):
        y = top + idx * (bar_h + gap)
        _textbox(slide, left, y, Inches(1.2), bar_h, str(label).title(), 10, RGBColor(69, 75, 89), bold=True)
        bar_w = width * 0.68 * (float(value) / max_value)
        bar = slide.shapes.add_shape(1, left + Inches(1.35), y, bar_w, bar_h)
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(34, 92, 130)
        bar.line.color.rgb = RGBColor(34, 92, 130)
        _textbox(slide, left + Inches(1.45) + bar_w, y, Inches(1.2), bar_h, _compact_money(value), 9, RGBColor(31, 41, 55))


def _draw_table_panel(slide, chart, left, top, width, height):
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    labels = chart.get("labels") or []
    if not labels:
        _textbox(slide, left, top, width, Inches(0.4), "Tabular result", 12, RGBColor(71, 85, 105), bold=True)
        return
    _textbox(slide, left, top, width, Inches(0.36), "Tabular result", 10, RGBColor(71, 85, 105), bold=True)
    row_h = Inches(0.36)
    for idx, label in enumerate(labels[:8]):
        y = top + Inches(0.5) + idx * row_h
        _textbox(slide, left, y, width, row_h, f"- {label}", 11, RGBColor(31, 41, 55))


def _business_status(answer: SQLAgentAnswer, quality: QualityReport) -> dict[str, str]:
    if answer.system_error is not None:
        return {
            "label": "Audit required",
            "tone": "critical",
            "description": "The live agent chain did not produce a decision-ready pack.",
        }
    if answer.question_id == AUDIT_ID:
        return {
            "label": "Audit required",
            "tone": "critical",
            "description": "The source conflict blocks a leadership-ready trend until analyst review is complete.",
        }
    if quality.layer_b.verdict == "DISAGREEMENT":
        return {
            "label": "Use with context",
            "tone": "caution",
            "description": "Ready for leadership when the source caveat travels with the recommendation.",
        }
    if answer.interpretation_choices:
        return {
            "label": "Needs analyst review",
            "tone": "review",
            "description": "Multiple definitions exist; choose the reporting definition before packaging.",
        }
    return {
        "label": "Ready to use",
        "tone": "ready",
        "description": "The result is ready to package into the decision pack.",
    }


def _business_question(answer: SQLAgentAnswer) -> str:
    if answer.question_id == HERO_ID:
        return "Which asset class should we prioritise for next month's growth plan?"
    if answer.question_id == AUDIT_ID:
        return "Can we rely on the monthly summary for the GTV trend?"
    if answer.question_id == MTU_ID:
        return "How many monthly transacting users should we report?"
    return answer.question


def _metric_name(answer: SQLAgentAnswer) -> str:
    labels = {
        HERO_ID: "GTV by asset class",
        "q2_gtv_usd_oct_2025": "Total GTV in USD",
        MTU_ID: "Monthly transacting users",
        "q4_transaction_count_by_asset_oct_2025": "Transaction count by asset class",
        AUDIT_ID: "GTV month-on-month trend",
    }
    return labels.get(answer.question_id, answer.metric_name.replace("_", " ").title())


def _headline(answer: SQLAgentAnswer, rows: list[dict[str, Any]]) -> str:
    if answer.question_id == HERO_ID and rows:
        top = max(rows, key=lambda r: float(r.get("gtv_idr") or 0))
        return f"Crypto is the clearest October growth priority at {_compact_money(top['gtv_idr'])} completed GTV."
    if answer.question_id == AUDIT_ID:
        return "Do not package the quarter-end trend until the stale source conflict is resolved."
    if answer.question_id == MTU_ID and rows:
        return "MTU is packable only after the reporting definition is chosen."
    return f"{_metric_name(answer)} is available for review."


def _decision_prompt(answer: SQLAgentAnswer) -> str:
    if answer.question_id == HERO_ID:
        return "Build a leadership pack for next month's asset-class growth focus."
    if answer.question_id == AUDIT_ID:
        return "Do not present the trend until the stale monthly summary conflict is resolved."
    return "Confirm the metric definition before building executive materials."


def _recommendation(answer: SQLAgentAnswer, status: dict[str, str]) -> str:
    if answer.question_id == HERO_ID:
        return "Lead with crypto as the priority growth lane, with gold and GSS as the secondary focus areas. Package the view as completed-transaction performance, not Ops activity."
    if answer.question_id == AUDIT_ID:
        return "Pause the leadership pack and route the trend to analyst review; the December monthly-summary total is stale relative to the canonical daily mart."
    if answer.question_id == MTU_ID:
        return "Select the MTU definition before packaging: AUM-defined, raw completed traders, or product-analytics activity."
    return status["description"]


def _recommended_use(answer: SQLAgentAnswer, status: dict[str, str]) -> str:
    if answer.question_id == HERO_ID:
        return "Recommended use: one-slide leadership recommendation on completed trading activity, with the Ops caveat included in speaker notes."
    if answer.question_id == AUDIT_ID:
        return "Recommended use: analyst audit brief only; not ready for a leadership recommendation."
    return f"Recommended use: {status['description']}"


def _source_caveat(answer: SQLAgentAnswer, quality: QualityReport) -> str:
    if answer.question_id == HERO_ID:
        return "Ops reporting runs higher because it includes pending transactions. Treat this as a source-policy difference, not a broken SQL result."
    if answer.question_id == AUDIT_ID:
        return "The monthly business summary has a known stale December total, so the portal blocks slide-ready use."
    if quality.layer_b.hypothesis is not None:
        return quality.layer_b.hypothesis.proposal
    return quality.layer_b.hypothesis_absence_note or "No material source caveat."


def _confidence(
    answer: SQLAgentAnswer, quality: QualityReport, status: dict[str, str]
) -> dict[str, str]:
    tp = quality.layer_c.trust_profile
    return {
        "business": status["label"],
        "correctness": tp.dimensions.correctness,
        "sourceReliability": tp.dimensions.source_reliability,
        "ambiguity": tp.dimensions.ambiguity,
        "overallEvidence": tp.overall,
    }


def _chart_spec(answer: SQLAgentAnswer, rows: list[dict[str, Any]]) -> dict[str, Any]:
    if not rows:
        return {"type": "table", "labels": [], "values": []}
    first = rows[0]
    if "asset_class" in first:
        value_key = "gtv_idr" if "gtv_idr" in first else "transaction_count"
        return {
            "type": "bar",
            "xKey": "asset_class",
            "yKey": value_key,
            "labels": [r["asset_class"] for r in rows],
            "values": [float(r[value_key]) for r in rows],
            "unit": "IDR" if "gtv" in value_key else "transactions",
        }
    if "month" in first:
        return {
            "type": "line",
            "xKey": "month",
            "yKey": "gtv_idr",
            "labels": [r["month"] for r in rows],
            "values": [float(r.get("gtv_idr") or 0) for r in rows],
            "unit": "IDR",
        }
    return {"type": "table", "labels": list(first.keys()), "values": []}


def _source_comparison(quality: QualityReport) -> list[dict[str, Any]]:
    out = []
    for finding in quality.layer_b.cross_source_findings:
        out.append(
            {
                "source": finding.source,
                "value": finding.value,
                "deltaVsPrimary": finding.delta_vs_primary,
                "notes": finding.notes,
            }
        )
    return out


def _executive_summary(
    answer: SQLAgentAnswer, quality: QualityReport, status: dict[str, str]
) -> list[str]:
    return [
        _headline(answer, _normalise_rows(answer.result_rows)),
        _recommendation(answer, status),
        _source_caveat(answer, quality),
        f"Decision-pack status: {status['label']} - {status['description']}",
    ]


def _workflow_steps(answer: SQLAgentAnswer, status: dict[str, str]) -> list[dict[str, str]]:
    package_state = "current" if status["label"] != "Audit required" else "blocked"
    review_state = "available" if status["label"] != "Audit required" else "current"
    return [
        {"label": "Ask", "state": "done", "detail": "Business question captured"},
        {"label": "Assess", "state": "done", "detail": status["label"]},
        {"label": "Package", "state": package_state, "detail": "Slide, CSV, email"},
        {"label": "Review", "state": review_state, "detail": "Evidence room"},
    ]


def _usage_boundaries(answer: SQLAgentAnswer, status: dict[str, str]) -> tuple[str, str]:
    if answer.question_id == HERO_ID:
        return (
            "Leadership recommendation on completed trading performance.",
            "Direct Ops-dashboard comparisons unless pending transactions are called out.",
        )
    if answer.question_id == AUDIT_ID:
        return (
            "Analyst audit discussion and source-owner follow-up.",
            "Board or leadership trend slides until December source quality is resolved.",
        )
    if answer.question_id == MTU_ID:
        return (
            "Definition-setting conversation before reporting MTU externally.",
            "A single headline MTU number without naming the chosen definition.",
        )
    return (status["description"], "External materials until the caveat is reviewed.")


def _chart_insight(answer: SQLAgentAnswer, rows: list[dict[str, Any]]) -> str:
    if answer.question_id == HERO_ID and rows:
        top = max(rows, key=lambda r: float(r.get("gtv_idr") or 0))
        total = sum(float(r.get("gtv_idr") or 0) for r in rows)
        share = float(top.get("gtv_idr") or 0) / total * 100 if total else 0
        return f"{str(top['asset_class']).title()} contributes {share:.0f}% of completed GTV in the selected period."
    if answer.question_id == AUDIT_ID:
        return "The trend shape is visible, but source disagreement blocks packaging."
    if answer.question_id == MTU_ID:
        return "Different measurement systems produce different MTU definitions."
    return "Chart is generated from executed SQL rows."


def _pack_title(answer: SQLAgentAnswer) -> str:
    if answer.question_id == HERO_ID:
        return "Asset-class growth priority pack"
    if answer.question_id == AUDIT_ID:
        return "Trend audit brief"
    if answer.question_id == MTU_ID:
        return "MTU definition decision pack"
    return f"{_metric_name(answer)} decision pack"


def _audit_payload(answer: SQLAgentAnswer, quality: QualityReport) -> dict[str, Any]:
    if answer.question_id != AUDIT_ID:
        return {"required": False, "reason": None, "nextActions": []}
    return {
        "required": True,
        "reason": "December monthly-summary totals conflict with the canonical completed-transaction mart.",
        "whatHappened": [
            "The SQL chain produced a canonical daily-mart trend.",
            "Layer B found the monthly business summary diverges beyond threshold.",
            "The portal blocks this result from becoming an executive slide until the stale source is resolved.",
        ],
        "nextActions": [
            "Ask the data owner to confirm whether the December Total row should be refreshed.",
            "Use fct_trading_daily for completed-transaction reporting while the conflict is open.",
            "Attach the analyst evidence appendix when discussing the discrepancy.",
        ],
        "unresolvedQuestions": quality.layer_c.unresolved_questions,
    }


def _normalise_rows(rows: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return [{str(k): v for k, v in row.items()} for row in rows]


def _compact_money(value: Any) -> str:
    number = float(value)
    if abs(number) >= 1_000_000_000:
        return f"{number / 1_000_000_000:.1f}B IDR"
    if abs(number) >= 1_000_000:
        return f"{number / 1_000_000:.1f}M"
    return f"{number:,.0f}"
